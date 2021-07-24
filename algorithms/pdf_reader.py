import PyPDF2
from pptx import Presentation
import docx
import re


def docx2string(filename):   #check later, may not get all text.
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return fullText

def pptx2string(filename):
    prs = Presentation(filename)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs

def remove_page_number(page_number: int, string):
    page_num_str = f"{page_number}"
    return re.sub(re.escape(page_num_str), "", string)

def pdf2string(filename, starting_page=0, end_page=None, *page_funcs):
    with open(filename, 'rb') as pdfFileObj:
        pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
        full_list =[]
        
        if not end_page:
            end_page = pdfReader.getNumPages()
        for i in range(starting_page, end_page):
            pageObj = pdfReader.getPage(i)
            unclean_text = pageObj.extractText()
            cleaned_text = remove_page_number(i, unclean_text)
            if page_funcs:
                for func in page_funcs:
                    cleaned_text = func(cleaned_text)
            full_list.append(cleaned_text)
    return full_list

def remove_lists_from_(the_list):
    for idx, el in enumerate(the_list):
        if type(el) == list:
            remove_lists_from_(el)
            list_length = len(el)
            for item in range(list_length):
                the_list.insert(idx, el.pop())
        if [] in the_list:
            the_list.remove([])
    return the_list

def text_cleaner(unclean_str):
    cleaned_str = re.sub("\\n", "", unclean_str)
    return cleaned_str

def filepath2string(filepath):
    return r'{}'.format(filepath)


def lect2vec(file_path):
    filename = filepath2string(file_path)
    is_pdf = bool(re.search("\.pdf$", filename))
    is_docx = bool(re.search("\.docx$", filename))
    is_pptx = bool(re.search("\.pptx$", filename))
    if is_pdf:
        list_of_lect_pages = pdf2string(file_path)
    elif is_docx:
        print("docx in lect2vec")
        list_of_lect_pages = docx2string(file_path)
    elif is_pptx:
        print("pptx in lect2vec")
        list_of_lect_pages = pptx2string(file_path)
    else:
        print("Docx, pptx, pdf only \n Try another file.")
    one_list = remove_lists_from_(list_of_lect_pages)
    one_string = ' '.join(one_list)
    cleaned_lecture = text_cleaner(one_string)
    return cleaned_lecture